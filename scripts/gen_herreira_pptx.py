"""Gera apresentação PowerPoint sobre cultura empresarial para Herreira Joias."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

GOLD = RGBColor(0xB8, 0x86, 0x0B)
GOLD_DARK = RGBColor(0x8B, 0x69, 0x14)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x2C, 0x2C, 0x2C)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x2E, 0x84, 0x4A)
BG_DARK = RGBColor(0x1A, 0x1A, 0x1A)
BG_LIGHT = RGBColor(0xFA, 0xF8, 0xF5)
GOLD_LIGHT = RGBColor(0xF5, 0xE6, 0xC8)

OUT = "C:/Sandyboxclaude/herreira_cultura_empresarial.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_quote(slide, text, author, top=5.8):
    add_text(slide, 0.8, top, 11.5, 0.8, f'"{text}"', size=14, color=GOLD, bold=False, font_name="Georgia", align=PP_ALIGN.CENTER)
    add_text(slide, 0.8, top + 0.5, 11.5, 0.4, f"— {author}", size=11, color=GRAY, align=PP_ALIGN.CENTER)


def add_bullet(slide, left, top, width, items, size=15, color=BLACK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(len(items) * 0.45 + 0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0


def gold_line(slide, top):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(top), Inches(2), Emu(36000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()


# === SLIDE 1: CAPA ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_DARK)
add_text(s, 0.8, 1.5, 11.5, 1.0, "HERREIRA JOIAS", size=16, color=GOLD, bold=True, align=PP_ALIGN.CENTER, font_name="Calibri")
add_text(s, 0.8, 2.2, 11.5, 1.5, "Transformacao Cultural:\nO Pilar Invisivel do Crescimento", size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Georgia")
add_text(s, 0.8, 4.2, 11.5, 0.6, "Estrategias para alinhar cultura, pessoas e resultado", size=18, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)
add_quote(s, "A cultura come a estrategia no cafe da manha.", "Peter Drucker", top=5.5)
add_text(s, 0.8, 6.5, 11.5, 0.4, "Assessoria Brasil GEO  |  2026", size=12, color=GRAY, align=PP_ALIGN.CENTER)

# === SLIDE 2: POR QUE CULTURA IMPORTA ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_LIGHT)
add_text(s, 0.8, 0.5, 5, 0.3, "POR QUE CULTURA IMPORTA", size=11, color=GOLD, bold=True)
gold_line(s, 0.9)
add_text(s, 0.8, 1.2, 11, 0.8, "Cultura: o ativo que ninguem ve\nmas que define tudo.", size=32, color=BLACK, bold=True, font_name="Georgia")
add_bullet(s, 0.8, 2.5, 5.5, [
    "Empresas com cultura forte crescem 3x mais rapido que concorrentes",
    "70% dos CEOs consideram cultura mais importante que estrategia",
    "+20% de produtividade em organizacoes com valores claros",
    "Cultura reduz turnover em ate 40% — economia direta em RH",
], size=14, color=GRAY)
add_bullet(s, 7, 2.5, 5.5, [
    "R$ 12,1M de receita na Herreira (+26,8%)",
    "Mas prejuizo de R$ 273K no grupo",
    "Zero politicas de RH, zero avaliacao de desempenho",
    "O crescimento sem cultura e insustentavel",
], size=14, color=GRAY)
add_quote(s, "As pessoas certas nos lugares certos fazem as coisas certas.", "Jim Collins, Good to Great")

# === SLIDE 3: DADOS ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_text(s, 0.8, 0.5, 5, 0.3, "DADOS QUE COMPROVAM", size=11, color=GOLD, bold=True)
gold_line(s, 0.9)
add_text(s, 0.8, 1.2, 11, 0.7, "O custo de ignorar a cultura e mensuravel.", size=32, color=BLACK, bold=True, font_name="Georgia")
# Stats
for i, (val, lab) in enumerate([("30%", "mais chance de crescer\ncom cultura forte"), ("R$ 1,4M", "perdidos por ano\nna Herreira por ineficiencia"), ("40%", "reducao de turnover\ncom valores claros"), ("3x", "mais receita em empresas\ncom cultura solida")]):
    x = 0.8 + i * 3.1
    add_text(s, x, 2.5, 2.8, 0.8, val, size=42, color=GOLD, bold=True, align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text(s, x, 3.4, 2.8, 0.8, lab, size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_quote(s, "O que e medido, e gerenciado.", "Peter Drucker")

# === SLIDE 4: HISTÓRIA DA PATRICIA ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_LIGHT)
add_text(s, 0.8, 0.5, 5, 0.3, "A HISTORIA", size=11, color=GOLD, bold=True)
gold_line(s, 0.9)
add_text(s, 0.8, 1.2, 11, 0.7, "De montagem de semijoias\na R$ 12 milhoes em receita.", size=32, color=BLACK, bold=True, font_name="Georgia")
add_bullet(s, 0.8, 2.5, 5.5, [
    "2008 — Patricia e Alexandre Caramaschi fundam a Herreira em Goiania",
    "2011 — Investem R$ 1 milhao em fabrica propria de semijoias",
    "110 funcionarios, vendas para todos os estados + EUA e Europa",
    "10-15 mil pecas/mes no atacado, presenca em novelas da Globo",
    "Marcas: Herreira + Aulore + Vitesse (Grupo HAV)",
], size=14, color=GRAY)
add_bullet(s, 7, 2.5, 5.5, [
    "56 anos de tradicao no mercado joalheiro de Goiania",
    "Celebridades como Monique Alfradique vestem Herreira",
    "Crescimento de 26,8% em 2025 — o melhor ano da Herreira",
    "Mas o grupo esta em prejuizo: hora de profissionalizar",
], size=14, color=GRAY)
add_quote(s, "Os resultados sao obtidos pela exploracao de oportunidades, nao pela solucao de problemas.", "Peter Drucker")

# === SLIDE 5: DIAGNÓSTICO RH ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_text(s, 0.8, 0.5, 5, 0.3, "DIAGNOSTICO: RECURSOS HUMANOS", size=11, color=RED, bold=True)
gold_line(s, 0.9)
add_text(s, 0.8, 1.2, 11, 0.7, "Zero politicas. Zero avaliacoes.\nZero plano de carreira.", size=32, color=BLACK, bold=True, font_name="Georgia")
add_bullet(s, 0.8, 2.5, 5.5, [
    "0 programas de capacitacao ou cursos",
    "0 mecanismos de retencao de talentos",
    "0 pesquisa de clima organizacional",
    "0 avaliacao de desempenho estruturada",
    "0 plano de carreira para colaboradores",
    "0 desenvolvimento de liderancas",
], size=14, color=RED)
add_bullet(s, 7, 2.5, 5.5, [
    "Impacto: alto turnover, baixa produtividade",
    "Custo invisivel: contratacao + treinamento repetido",
    "Aulore: pessoal consome 38,5% da receita",
    "Sem comunicacao interna estruturada",
    "Sem mecanismo de valorizacao 'Prata da Casa'",
], size=14, color=GRAY)
add_quote(s, "Nao ha nada tao inutil quanto fazer com grande eficiencia algo que nao deveria ser feito.", "Peter Drucker")

# === SLIDE 6: GOVERNANÇA ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_LIGHT)
add_text(s, 0.8, 0.5, 5, 0.3, "DIAGNOSTICO: GOVERNANCA", size=11, color=RED, bold=True)
gold_line(s, 0.9)
add_text(s, 0.8, 1.2, 11, 0.7, "A unica coisa que funciona:\nseparacao patrimonial.", size=32, color=BLACK, bold=True, font_name="Georgia")
add_bullet(s, 0.8, 2.5, 5.5, [
    "Separacao empresa/socias: nota 10/10",
    "Impostos em dia: zero atrasados",
    "Financiamentos regulares e pagos",
], size=14, color=GREEN)
add_bullet(s, 7, 2.5, 5.5, [
    "Zero compliance, zero auditoria",
    "LGPD nao implementada (multa ate R$ 50M)",
    "Sem mapa de gestao de riscos",
    "Sem programa de sucessao familiar",
    "Codigo de conduta pouco divulgado (nota 4/10)",
], size=14, color=RED)
add_quote(s, "Gerenciar e fazer as coisas direito. Liderar e fazer as coisas certas.", "Peter Drucker")

# === SLIDE 7: PRO LABORE ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_text(s, 0.8, 0.5, 5, 0.3, "O PROBLEMA CENTRAL", size=11, color=RED, bold=True)
gold_line(s, 0.9)
add_text(s, 0.8, 1.2, 11, 0.7, "Socias retiram mais do que\na empresa consegue gerar.", size=32, color=BLACK, bold=True, font_name="Georgia")
for i, (unit, pay, sit) in enumerate([("Herreira", "-60,9%", "Controlavel"), ("Aulore", "-363%", "CRITICO"), ("Vitesse", "-230%", "CRITICO")]):
    x = 0.8 + i * 4
    color = GREEN if i == 0 else RED
    add_text(s, x, 2.8, 3.5, 0.5, unit, size=24, color=GOLD, bold=True, font_name="Georgia", align=PP_ALIGN.CENTER)
    add_text(s, x, 3.4, 3.5, 0.5, f"Payout: {pay}", size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, 3.9, 3.5, 0.4, sit, size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 4.6, 11, 0.6, "Resultado: descapitalizacao progressiva. Pro labore total das 3 socias: R$ 1,6M/ano.\nLucro operacional do grupo: R$ 1,5M. Prejuizo apos retiradas: -R$ 273K.", size=14, color=GRAY)
add_quote(s, "Lucro nao e um objetivo. E uma condicao necessaria para a sobrevivencia.", "Peter Drucker")

# === SLIDE 8: CUSTO DE NÃO TER CULTURA ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_DARK)
add_text(s, 0.8, 0.5, 5, 0.3, "O CUSTO DA INACAO", size=11, color=GOLD, bold=True)
add_text(s, 0.8, 1.2, 11, 0.7, "R$ 1,4 milhao por ano.\nEsse e o preco de nao mudar.", size=36, color=WHITE, bold=True, font_name="Georgia")
for i, (val, lab) in enumerate([("R$ 760K", "Custos de produto\nacima do benchmark"), ("R$ 366K", "Ineficiencia\nem pessoal"), ("R$ 360K", "Despesas financeiras\nrenegociaveis")]):
    x = 1.5 + i * 3.8
    add_text(s, x, 2.8, 3.2, 0.7, val, size=36, color=GOLD, bold=True, align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text(s, x, 3.6, 3.2, 0.8, lab, size=13, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)
add_quote(s, "Se voce quer algo novo, voce precisa parar de fazer algo velho.", "Peter Drucker", top=5.2)

# === SLIDE 9: 6 PILARES ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_LIGHT)
add_text(s, 0.8, 0.5, 5, 0.3, "O QUE MUDAR", size=11, color=GOLD, bold=True)
gold_line(s, 0.9)
add_text(s, 0.8, 1.2, 11, 0.7, "6 pilares para transformar\na cultura da Herreira.", size=32, color=BLACK, bold=True, font_name="Georgia")
pilares = [
    ("01", "Valores", "Definir e comunicar os valores que guiam decisoes"),
    ("02", "Rituais", "Reunioes mensais de resultado, celebracoes, feedbacks"),
    ("03", "Metricas", "KPIs por funcao, avaliacao de desempenho trimestral"),
    ("04", "Lideranca", "Desenvolver gestoras como exemplos vivos da cultura"),
    ("05", "Desenvolvimento", "Capacitacao continua, plano de carreira, mentoria"),
    ("06", "Reconhecimento", "Programa Prata da Casa, bonus por performance"),
]
for i, (num, title, desc) in enumerate(pilares):
    row = i // 3
    col = i % 3
    x = 0.8 + col * 4
    y = 2.5 + row * 2
    add_text(s, x, y, 0.5, 0.4, num, size=14, color=GOLD, bold=True)
    add_text(s, x + 0.6, y, 3, 0.4, title, size=18, color=BLACK, bold=True)
    add_text(s, x + 0.6, y + 0.4, 3, 0.5, desc, size=12, color=GRAY)

# === SLIDE 10: PLANO 90 DIAS ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, WHITE)
add_text(s, 0.8, 0.5, 5, 0.3, "PLANO DE ACAO", size=11, color=GOLD, bold=True)
gold_line(s, 0.9)
add_text(s, 0.8, 1.2, 11, 0.7, "90 dias para reverter o prejuizo\ne construir a base cultural.", size=32, color=BLACK, bold=True, font_name="Georgia")
fases = [
    ("Semana 1-4", "Diagnostico + financeiro", "Mapear custos, renegociar dividas, limitar pro labore"),
    ("Semana 2-6", "CRM + pos-venda", "Implementar CRM, NPS, reativar clientes inativos"),
    ("Semana 3-8", "Marketing digital", "Instagram Shopping, Reels, Google Ads, influenciadoras"),
    ("Semana 4-8", "Governanca basica", "Rituais mensais, compliance, LGPD, mapa de riscos"),
    ("Semana 6-10", "Pessoas e cultura", "Manual do colaborador, avaliacao, plano de carreira"),
    ("Semana 8-12", "Escala digital", "E-commerce, marketplace, personal branding Paty"),
]
for i, (when, what, how) in enumerate(fases):
    y = 2.4 + i * 0.7
    add_text(s, 0.8, y, 2, 0.4, when, size=12, color=GOLD, bold=True)
    add_text(s, 3, y, 3, 0.4, what, size=14, color=BLACK, bold=True)
    add_text(s, 6.5, y, 6, 0.4, how, size=12, color=GRAY)
add_quote(s, "Planejamento de longo prazo nao lida com decisoes futuras, mas com o futuro das decisoes atuais.", "Peter Drucker", top=6.5)

# === SLIDE 11: BENCHMARK VIVARA ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_LIGHT)
add_text(s, 0.8, 0.5, 5, 0.3, "BENCHMARK", size=11, color=GOLD, bold=True)
gold_line(s, 0.9)
add_text(s, 0.8, 1.2, 11, 0.7, "Herreira vs Vivara: onde\na cultura faz a diferenca.", size=32, color=BLACK, bold=True, font_name="Georgia")
metricas = [
    ("Crescimento", "+26,8%", "+16,2%", "Herreira GANHA"),
    ("Margem bruta", "52,6%", "70-74%", "Gap de 18pp"),
    ("EBITDA", "12,7%", "25,5%", "Gap de 13pp"),
    ("Margem liquida", "-1,3%", "25,4%", "CRITICO"),
    ("E-commerce", "Incipiente", "Robusto", "Gap critico"),
]
for i, (met, her, viv, gap) in enumerate(metricas):
    y = 2.6 + i * 0.65
    add_text(s, 0.8, y, 2.5, 0.4, met, size=13, color=BLACK, bold=True)
    add_text(s, 3.5, y, 2.5, 0.4, her, size=13, color=GOLD)
    add_text(s, 6, y, 2.5, 0.4, viv, size=13, color=GRAY)
    add_text(s, 8.5, y, 3, 0.4, gap, size=12, color=RED if "CRITICO" in gap or "Gap" in gap else GREEN)
add_quote(s, "Voce nao pode gerenciar o que nao pode medir.", "Peter Drucker")

# === SLIDE 12: CASE INSPIRADOR ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_DARK)
add_text(s, 0.8, 0.5, 5, 0.3, "INSPIRACAO", size=11, color=GOLD, bold=True)
add_text(s, 0.8, 1.5, 11, 1.0, "A Vivara era uma joalheria familiar.\nHoje esta na Nasdaq.", size=36, color=WHITE, bold=True, font_name="Georgia", align=PP_ALIGN.CENTER)
add_text(s, 1.5, 3.2, 10, 1.5, "A Vivara foi fundada como joalheria familiar em 1962. Profissionalizou a gestao, implementou governanca corporativa, abriu capital em 2019 (VIVA3). Hoje fatura R$ 3,8 bilhoes com 400+ lojas. O caminho da Herreira nao precisa ser o IPO — mas a profissionalizacao da gestao e o primeiro passo para qualquer empresa familiar que quer crescer de forma sustentavel.", size=16, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)
add_quote(s, "Primeiro quem, depois o que. Coloque as pessoas certas no onibus.", "Jim Collins, Good to Great", top=5.5)

# === SLIDE 13: PROJEÇÃO ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_LIGHT)
add_text(s, 0.8, 0.5, 5, 0.3, "PROJECAO FINANCEIRA", size=11, color=GOLD, bold=True)
gold_line(s, 0.9)
add_text(s, 0.8, 1.2, 11, 0.7, "De prejuizo a lucro em 12 meses.\nSem investimento adicional.", size=32, color=BLACK, bold=True, font_name="Georgia")
for i, (val, lab, col) in enumerate([("-R$ 273K", "Resultado atual\n(2025)", RED), ("+R$ 1,4M", "Melhorias\nidentificadas", GOLD), ("+R$ 1,1M", "Projecao\n(2026)", GREEN)]):
    x = 1 + i * 4
    add_text(s, x, 2.8, 3.5, 0.7, val, size=40, color=col, bold=True, align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text(s, x, 3.6, 3.5, 0.6, lab, size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 4.8, 11, 0.5, "Investimento necessario: R$ 0 (gestao operacional) + R$ 3K-8K (governanca) + R$ 3K-8K/mes (marketing digital)", size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_quote(s, "A melhor maneira de prever o futuro e cria-lo.", "Peter Drucker")

# === SLIDE 14: FRASE FINAL ===
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, BG_DARK)
add_text(s, 0.8, 2.0, 11.5, 1.5, "Cultura nao e o que voce\nescreve na parede.\nE o que acontece quando\nninguem esta olhando.", size=38, color=WHITE, bold=True, font_name="Georgia", align=PP_ALIGN.CENTER)
add_text(s, 0.8, 4.5, 11.5, 0.5, "— Adaptado de Peter Drucker", size=16, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 5.8, 11.5, 0.5, "HERREIRA JOIAS  |  O proximo capitulo comeca agora.", size=14, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)
add_text(s, 0.8, 6.5, 11.5, 0.4, "Assessoria por Alexandre Caramaschi  |  CEO da Brasil GEO  |  alexandrecaramaschi.com", size=11, color=GRAY, align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f"Apresentacao salva em: {OUT}")
print(f"Total de slides: {len(prs.slides)}")
